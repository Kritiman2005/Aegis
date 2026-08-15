import os
import uuid
import logging
from typing import List, Dict, Any
from pathlib import Path

# Extractors
import fitz  # PyMuPDF
from pptx import Presentation
import easyocr

# Vector DB & Embeddings
from qdrant_client import QdrantClient, models
from fastembed import TextEmbedding, SparseTextEmbedding
from sentence_transformers import CrossEncoder

logger = logging.getLogger(__name__)

# ─── Configuration ────────────────────────────────────────────────────────────

BASE_DIR = Path(__file__).resolve().parent.parent.parent.parent
data_dir = os.environ.get("AEGIS_DATA_DIR")
if data_dir:
    QDRANT_DB_DIR = Path(data_dir) / "qdrant_db"
else:
    QDRANT_DB_DIR = BASE_DIR / "qdrant_db"

COLLECTION_NAME = "aegis_documents"

# Initialize Qdrant persistent client on the event loop thread
_qdrant_client = None

def init_qdrant():
    """Called on FastAPI startup to bind QdrantClient to the main event loop thread."""
    global _qdrant_client
    if _qdrant_client is None:
        logger.info("Initializing QdrantClient on Uvicorn event loop thread...")
        _qdrant_client = QdrantClient(path=str(QDRANT_DB_DIR))
        # Ensure Collection Exists
        if not _qdrant_client.collection_exists(COLLECTION_NAME):
            _qdrant_client.create_collection(
                collection_name=COLLECTION_NAME,
                vectors_config={
                    "text-dense": models.VectorParams(
                        size=384,  # BAAI/bge-small-en-v1.5 output size
                        distance=models.Distance.COSINE
                    )
                },
                sparse_vectors_config={
                    "text-sparse": models.SparseVectorParams(
                        modifier=models.Modifier.IDF
                    )
                }
            )

def get_qdrant_client():
    if _qdrant_client is None:
        init_qdrant()
    return _qdrant_client

# Initialize FastEmbed models (Downloaded automatically on first run)
logger.info("Initializing Dense & Sparse Embedding Models...")
dense_model = TextEmbedding("BAAI/bge-small-en-v1.5")
sparse_model = SparseTextEmbedding("Qdrant/bm25")

# Initialize Reranker
_reranker = None
def get_reranker():
    global _reranker
    if _reranker is None:
        logger.info("Initializing CrossEncoder Reranker...")
        _reranker = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2')
    return _reranker

# Initialize EasyOCR reader (Lazy load)
_ocr_reader = None

def get_ocr_reader():
    global _ocr_reader
    if _ocr_reader is None:
        logger.info("Initializing EasyOCR reader (this might take a moment)...")
        _ocr_reader = easyocr.Reader(['en'], gpu=False)
    return _ocr_reader


# ─── Text Extraction ──────────────────────────────────────────────────────────

def extract_text(file_path: str, file_type: str) -> str:
    ext = file_type.lower()
    try:
        if ext == 'pdf':
            doc = fitz.open(file_path)
            text_content = ""
            for page in doc:
                text_content += page.get_text() + "\n"
            return text_content
            
        elif ext in ['ppt', 'pptx']:
            prs = Presentation(file_path)
            text_content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
            return text_content
            
        elif ext in ['txt', 'md', 'csv']:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
                
        elif ext in ['png', 'jpg', 'jpeg']:
            reader = get_ocr_reader()
            results = reader.readtext(file_path)
            return " ".join([res[1] for res in results])
            
        else:
            logger.warning(f"Unsupported file type for extraction: {ext}")
            return ""
    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        raise e


# ─── Chunking ─────────────────────────────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = 300, overlap: int = 50) -> List[str]:
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i:i + chunk_size])
        chunks.append(chunk)
        i += (chunk_size - overlap)
    return chunks


# ─── Ingestion ────────────────────────────────────────────────────────────────

def ingest_document(document_id: int, file_path: str, file_type: str, filename: str):
    logger.info(f"Ingesting document {document_id}: {filename}")
    
    raw_text = extract_text(file_path, file_type)
    if not raw_text.strip():
        logger.warning(f"Document {document_id} resulted in empty text.")
        return
        
    chunks = chunk_text(raw_text)
    logger.info(f"Generated {len(chunks)} chunks for document {document_id}.")
    
    # Generate Dense and Sparse Vectors
    dense_vecs = list(dense_model.embed(chunks))
    sparse_vecs = list(sparse_model.embed(chunks))
    
    points = []
    for i, chunk in enumerate(chunks):
        # fastembed sparse returns SparseEmbedding object
        sv = sparse_vecs[i]
        
        point = models.PointStruct(
            id=str(uuid.uuid4()),
            vector={
                "text-dense": dense_vecs[i].tolist(),
                "text-sparse": models.SparseVector(
                    indices=sv.indices.tolist(),
                    values=sv.values.tolist()
                )
            },
            payload={
                "document_id": document_id,
                "chunk_index": i,
                "filename": filename,
                "content": chunk
            }
        )
        points.append(point)
        
    client = get_qdrant_client()
    client.upsert(
        collection_name=COLLECTION_NAME,
        points=points
    )
    logger.info(f"Successfully ingested document {document_id} into Qdrant.")


# ─── Advanced Hybrid Retrieval & Reranking ────────────────────────────────────

def hybrid_search(query: str, conversation_id: str, top_k: int = 5) -> List[Dict[str, Any]]:
    """
    1. Qdrant Native Hybrid Search (Fusion)
    2. Filter by Conversation ID
    3. Rerank via CrossEncoder to get Top K
    """
    logger.info(f"Running Qdrant hybrid search for query: {query}")
    
    # We must first fetch the user's document IDs for this conversation
    from app.db.database import SessionLocal
    from app.db.models import UserDocument
    db = SessionLocal()
    docs = db.query(UserDocument).filter(UserDocument.conversation_id == conversation_id).all()
    db.close()
    
    valid_doc_ids = [d.id for d in docs]
    if not valid_doc_ids:
        return []
    
    query_dense = list(dense_model.embed([query]))[0]
    query_sparse = list(sparse_model.embed([query]))[0]
    
    doc_filter = models.Filter(
        must=[
            models.FieldCondition(
                key="document_id",
                match=models.MatchAny(any=valid_doc_ids)
            )
        ]
    )

    # Query Qdrant with Reciprocal Rank Fusion (RRF) implicitly by querying both
    # Qdrant's query_points automatically fuses multiple prefetches
    results = get_qdrant_client().query_points(
        collection_name=COLLECTION_NAME,
        prefetch=[
            models.Prefetch(
                query=query_dense.tolist(),
                using="text-dense",
                limit=15,
                filter=doc_filter
            ),
            models.Prefetch(
                query=models.SparseVector(
                    indices=query_sparse.indices.tolist(),
                    values=query_sparse.values.tolist()
                ),
                using="text-sparse",
                limit=15,
                filter=doc_filter
            )
        ],
        query=models.FusionQuery(fusion=models.Fusion.RRF),
        limit=15
    )
    
    unique_chunks = []
    for point in results.points:
        unique_chunks.append({
            "id": point.id,
            "content": point.payload.get("content", ""),
            "document_id": point.payload.get("document_id"),
            "filename": point.payload.get("filename"),
            "fusion_score": point.score
        })
        
    if not unique_chunks:
        return []
        
    # Rerank
    logger.info(f"Reranking {len(unique_chunks)} fused chunks...")
    reranker = get_reranker()
    pairs = [[query, chunk["content"]] for chunk in unique_chunks]
    scores = reranker.predict(pairs)
    
    filtered_chunks = []
    for i, chunk in enumerate(unique_chunks):
        score = float(scores[i])
        chunk["rerank_score"] = score
        filtered_chunks.append(chunk)
            
    filtered_chunks.sort(key=lambda x: x["rerank_score"], reverse=True)
    return filtered_chunks[:top_k]
